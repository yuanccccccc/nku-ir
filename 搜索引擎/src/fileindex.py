import os
import csv
import requests
from urllib.parse import urlparse
from pptx import Presentation
import pandas as pd
import pdfplumber
from elasticsearch import Elasticsearch
import win32com.client  # 用于将 .doc 文件转换为 .docx 文件

# 初始化 Elasticsearch 客户端
es = Elasticsearch([{"host": "localhost", "port": 9200, "scheme": "http"}])

# 定义文件下载目录
DOWNLOAD_FOLDER = "D:\\test\\ir\\project\\files"
os.makedirs(DOWNLOAD_FOLDER, exist_ok=True)

# CSV 文件路径
CSV_PATH = "D:\\test\\ir\\project\\webpages.csv"

# 索引名称
INDEX_NAME = "fileindex"

# 索引映射设置
INDEX_MAPPING = {
    "mappings": {
        "properties": {
            "url": {"type": "keyword"},
            "content": {"type": "text", "analyzer": "ik_smart"},
        }
    }
}


def check_create_index(index_name):
    """
    检查索引是否存在，如果不存在则创建索引
    """
    if es.indices.exists(index=index_name):
        es.indices.delete(index=index_name)
    es.indices.create(index=index_name, body=INDEX_MAPPING)


def download_file(url, download_folder):
    """
    下载文件并保存到指定的目录
    """
    try:
        response = requests.get(url)
        response.raise_for_status()
        filename = os.path.basename(urlparse(url).path)
        file_path = os.path.join(download_folder, filename)
        with open(file_path, "wb") as f:
            f.write(response.content)
        return file_path, filename
    except requests.exceptions.RequestException as e:
        print(f"Error downloading {url}: {e}")
        return None, None


def extract_text(file_path, file_extension):
    """
    根据文件扩展名提取文件内容
    """
    text = ""
    try:
        if file_extension == "pdf":
            text = extract_text_from_pdf(file_path)
        elif file_extension == "doc":
            text = extract_text_from_doc(file_path)
        elif file_extension in ["ppt", "pptx"]:
            text = extract_text_from_ppt(file_path)
        elif file_extension in ["xls", "xlsx"]:
            text = extract_text_from_excel(file_path)
        print(f"Extracted text {text} from {file_path}")
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
    return text


def extract_text_from_pdf(pdf_path):
    """
    Extract all text content from a PDF file using PyPDF2.

    :param pdf_path: Path to the PDF file
    :return: A string containing all text in the document
    """
    try:
        with open(pdf_path, "rb") as file:
            with pdfplumber.open(file) as pdf:
                text = ""
                for page in pdf.pages:
                    text += page.extract_text()
        text = text.replace("\n", "").replace(" ", "")
        return text
    except Exception:
        return ""


def extract_text_from_doc(doc_path):
    """
    Extract all text content from a .doc file using win32com.

    :param doc_path: Path to the .doc file
    :return: A string containing all text in the document
    """
    try:
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False  # Run in the background
        word.DisplayAlerts = False  # Suppress alerts

        # Open the document with forced editing mode
        doc = word.Documents.Open(doc_path, ReadOnly=False, AddToRecentFiles=False)

        # Activate editing mode if in Protected View
        if doc.ReadOnly:
            doc.ReadOnlyRecommended = False

        text = doc.Content.Text  # Extract the content text

        # Close the document and quit Word
        doc.Close(False)  # Close without saving
        word.Quit()

        return text.replace("\n", "").replace(" ", "")
    except Exception:
        return ""


def extract_text_from_ppt(ppt_path):
    """
    Extract all text content from a PPT file using python-pptx.

    :param ppt_path: Path to the PPT file
    :return: A string containing all text in the document
    """
    try:
        prs = Presentation(ppt_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
        text = text.replace("\n", "").replace(" ", "")
        return text
    except Exception as e:
        return ""


def extract_text_from_excel(file_path):
    """
    提取 Excel 文件中的文本
    """
    try:
        df = pd.read_excel(file_path)
        return df.to_string().replace("\n", "").replace(" ", "")
    except Exception as e:
        print(f"Error extracting text from Excel {file_path}: {e}")
        return ""


def index_document(url, content):
    """
    将文件的 URL 和内容索引到 Elasticsearch 中
    """
    doc = {"url": url, "content": content}
    try:
        es.index(index=INDEX_NAME, document=doc)
        print(f"Document indexed: {url}")
    except Exception as e:
        print(f"Error indexing document {url}: {e}")


def process_csv(csv_path):
    """
    处理 CSV 文件，下载文件并索引文件内容
    """
    with open(csv_path, newline="", encoding="utf-8") as csvfile:
        csvreader = csv.DictReader(csvfile)
        i = 0
        for row in csvreader:
            url = row["URL"]
            # 判断文件扩展名，是否为支持的文件类型
            if is_supported_file(url):
                print(f"Processing: {url}")
                file_path, filename = download_file(url, DOWNLOAD_FOLDER)
                if file_path:
                    file_extension = filename.split(".")[-1].lower()
                    content = extract_text(file_path, file_extension).replace(" ", "")
                    index_document(url, content)


def is_supported_file(url):
    """
    判断 URL 是否为支持的文件类型
    """
    return url.lower().endswith((".pdf", ".doc", ".docx", ".pptx", ".xls", ".xlsx"))


# 检查并创建索引
check_create_index(INDEX_NAME)
# 处理 CSV 文件
process_csv(CSV_PATH)
